#!/usr/bin/env python3
"""
Render an approximate PNG of each slide in a .pptx, without PowerPoint.

Useful for checking a conversion at a glance: it draws the slide's background
picture and then every text box at its real position, size, colour and weight.
Line breaking and font metrics are approximations, so treat it as a proof sheet
rather than a pixel-accurate render.

USAGE
    python preview_pptx.py output.pptx previews/
    python preview_pptx.py output.pptx previews/ --pages 2,3 --width 1600
"""
import argparse
import io
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

EMU_PER_INCH = 914400.0


def load_font(path, size_px, cache={}):
    key = (path, size_px)
    if key not in cache:
        try:
            cache[key] = ImageFont.truetype(path, max(size_px, 1))
        except Exception:
            cache[key] = ImageFont.load_default()
    return cache[key]


def render_slide(slide, slide_w_emu, slide_h_emu, width_px, regular, bold):
    px_per_emu = width_px / float(slide_w_emu)
    height_px = int(round(slide_h_emu * px_per_emu))
    canvas = Image.new("RGB", (width_px, height_px), "white")

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
            box = (
                int(shape.left * px_per_emu), int(shape.top * px_per_emu),
                max(int(shape.width * px_per_emu), 1), max(int(shape.height * px_per_emu), 1),
            )
            canvas.paste(picture.resize((box[2], box[3]), Image.LANCZOS), (box[0], box[1]))

    draw = ImageDraw.Draw(canvas)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        left, top = shape.left * px_per_emu, shape.top * px_per_emu
        box_w = shape.width * px_per_emu
        cursor = top
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs)
            if not text:
                continue
            first = paragraph.runs[0].font
            size_pt = first.size.pt if first.size else 12.0
            size_px = int(round(size_pt / 72.0 * EMU_PER_INCH * px_per_emu))

            # Draw run by run: a line may mix colours and weights mid-sentence.
            styled = []
            for run in paragraph.runs:
                font = load_font(bold if run.font.bold else regular, size_px)
                try:
                    colour = tuple(int(v) for v in bytes.fromhex(str(run.font.color.rgb)))
                except Exception:
                    colour = (0, 0, 0)
                styled.append((run.text, font, colour))

            text_w = sum(draw.textlength(t, font=f) for t, f, _ in styled)
            x = left
            if paragraph.alignment == PP_ALIGN.CENTER:
                x = left + (box_w - text_w) / 2
            elif paragraph.alignment == PP_ALIGN.RIGHT:
                x = left + box_w - text_w
            for chunk, font, colour in styled:
                draw.text((x, cursor), chunk, font=font, fill=colour)
                x += draw.textlength(chunk, font=font)
            cursor += size_px * 1.22
    return canvas


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path")
    ap.add_argument("out_dir")
    ap.add_argument("--width", type=int, default=1600, help="Preview width in pixels")
    ap.add_argument("--pages", help='Slides to render, 1-based, e.g. "2,3"')
    ap.add_argument("--font-file", default=r"C:\Windows\Fonts\arial.ttf")
    ap.add_argument("--bold-font-file", default=r"C:\Windows\Fonts\arialbd.ttf")
    args = ap.parse_args()

    wanted = None
    if args.pages:
        wanted = {int(p) - 1 for p in args.pages.replace(" ", "").split(",") if p}

    prs = Presentation(args.pptx_path)
    os.makedirs(args.out_dir, exist_ok=True)
    for index, slide in enumerate(prs.slides):
        if wanted is not None and index not in wanted:
            continue
        image = render_slide(slide, prs.slide_width, prs.slide_height,
                             args.width, args.font_file, args.bold_font_file)
        path = os.path.join(args.out_dir, f"slide{index + 1}.png")
        image.save(path)
        print(f"wrote {path}")


if __name__ == "__main__":
    main()
