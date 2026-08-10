#!/usr/bin/env python3
"""
OCR a scanned PDF slide deck into a fully editable PowerPoint (.pptx).

Every page becomes one slide. Detected text is turned into real, editable
PowerPoint text boxes at its original position, size, colour and weight. The
original page is kept as the slide background with the captured text cleanly
removed, so nothing is printed twice.

CORE INVARIANT
    Text is erased from the background ONLY if it was successfully captured as
    an editable text box. Anything OCR could not read stays visible as pixels.
    No text can silently disappear from a slide.

PIPELINE
    1. Render each page well above its native resolution.
    2. Detect words in both polarities (normal + inverted) so light-on-dark
       text -- navy table headers, coloured banners -- is found too.
    3. Cluster words into rows geometrically, then split each row at wide
       horizontal gaps so table cells stay separate.
    4. Re-OCR every uncertain line as an isolated, upscaled, auto-polarity
       crop, and keep the rewrite only when it measurably improves.
    5. Drop what is still unreadable; measure colour/weight/size of the rest
       from its own pixels.
    6. Erase the captured text: flat fill where the background is uniform,
       inpainting where it is not.
    7. Emit one text box per paragraph, with detected alignment.

USAGE
    python ocr_pdf_to_pptx.py input.pdf output.pptx
    python ocr_pdf_to_pptx.py input.pdf output.pptx --lang ind+eng --dpi 400
    python ocr_pdf_to_pptx.py input.pdf output.pptx --background original
"""
import argparse
import io
import os
import shutil
import statistics
import sys

import cv2
import numpy as np
import pymupdf
import pytesseract
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Tesseract discovery (the Windows installer puts it outside PATH)
# --------------------------------------------------------------------------
TESSERACT_CANDIDATES = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe"),
]


def locate_tesseract(explicit=None):
    if explicit and os.path.isfile(explicit):
        return explicit
    found = shutil.which("tesseract")
    if found:
        return found
    for cand in TESSERACT_CANDIDATES:
        if os.path.isfile(cand):
            return cand
    raise SystemExit(
        "Tesseract not found. Install it (see README.txt) or pass --tesseract "
        "with the full path to tesseract.exe."
    )


def render_page(page, dpi):
    zoom = dpi / 72.0
    pix = page.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
    return np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3).copy()


# --------------------------------------------------------------------------
# Word detection
# --------------------------------------------------------------------------
def detect_words(rgb, lang):
    """Find candidate words in both polarities.

    Deliberately keeps low-confidence hits. Confidence filtering happens AFTER
    the per-line refinement pass -- filtering here is what made v1 lose the
    navy table header, which the page pass reads as "ed" (conf 18) but which a
    cropped re-OCR reads perfectly as "Inisial".
    """
    words = []
    for polarity, image in (("normal", rgb), ("inverted", 255 - rgb)):
        data = pytesseract.image_to_data(
            Image.fromarray(image), lang=lang, output_type=pytesseract.Output.DICT
        )
        for i in range(len(data["text"])):
            text = data["text"][i].strip()
            if not text:
                continue
            try:
                conf = float(data["conf"][i])
            except (TypeError, ValueError):
                continue
            if conf < 0:
                continue
            words.append(
                {
                    "text": text,
                    "conf": conf,
                    "x0": data["left"][i],
                    "y0": data["top"][i],
                    "x1": data["left"][i] + data["width"][i],
                    "y1": data["top"][i] + data["height"][i],
                    "polarity": polarity,
                }
            )
    return dedupe_boxes(words)


def iou(a, b):
    ix0, iy0 = max(a["x0"], b["x0"]), max(a["y0"], b["y0"])
    ix1, iy1 = min(a["x1"], b["x1"]), min(a["y1"], b["y1"])
    inter = max(0, ix1 - ix0) * max(0, iy1 - iy0)
    if inter == 0:
        return 0.0
    area_a = (a["x1"] - a["x0"]) * (a["y1"] - a["y0"])
    area_b = (b["x1"] - b["x0"]) * (b["y1"] - b["y0"])
    return inter / float(area_a + area_b - inter)


def containment(a, b):
    """Fraction of a's area lying inside b.

    Catches fragments that IoU misses: a stray "Laik" re-detected inside an
    already-captured "Laik Kerja" scores low IoU but is fully contained.
    """
    inter = (max(0, min(a["x1"], b["x1"]) - max(a["x0"], b["x0"]))
             * max(0, min(a["y1"], b["y1"]) - max(a["y0"], b["y0"])))
    area = max((a["x1"] - a["x0"]) * (a["y1"] - a["y0"]), 1)
    return inter / float(area)


def dedupe_boxes(words, thresh=0.5):
    """Drop near-duplicate detections, keeping the higher-confidence one.

    Both polarity passes normally find dark-on-light text, so without this
    every ordinary word would be detected twice.
    """
    kept = []
    for w in sorted(words, key=lambda w: -w["conf"]):
        if not any(iou(w, k) > thresh for k in kept):
            kept.append(w)
    return kept


# --------------------------------------------------------------------------
# Rows -> lines
# --------------------------------------------------------------------------
def same_row(a, b):
    """Do two boxes sit on the same baseline?

    Requires comparable heights and near-identical vertical centres. A plain
    overlap test is not enough: a tall heading's bbox fully contains the small
    body line beneath it, which welded headings into body text in v2.
    """
    ha, hb = a["y1"] - a["y0"], b["y1"] - b["y0"]
    if max(ha, hb) / max(min(ha, hb), 1) > 2.2:
        return False
    inter = max(0, min(a["y1"], b["y1"]) - max(a["y0"], b["y0"]))
    if inter / float(max(min(ha, hb), 1)) < 0.5:
        return False
    ca, cb = (a["y0"] + a["y1"]) / 2.0, (b["y0"] + b["y1"]) / 2.0
    return abs(ca - cb) <= 0.45 * min(ha, hb)


def assemble_lines(words, gap_factor=1.0):
    """Cluster words into rows, then split rows at wide horizontal gaps.

    Grouping must NOT use Tesseract's block/par/line numbers: those are
    per-pass, so after merging two polarity passes the same visual line carries
    two different keys and gets torn into fragments. Splitting on gaps keeps
    table cells apart, since Tesseract reports a whole table row as one "line".
    """
    rows = []
    for w in sorted(words, key=lambda w: (w["y0"], w["x0"])):
        candidates = [r for r in rows if same_row(w, r)]
        if candidates:
            centre = (w["y0"] + w["y1"]) / 2.0
            row = min(candidates, key=lambda r: abs((r["y0"] + r["y1"]) / 2.0 - centre))
            row["words"].append(w)
            row["y0"] = min(row["y0"], w["y0"])
            row["y1"] = max(row["y1"], w["y1"])
        else:
            rows.append({"words": [w], "y0": w["y0"], "y1": w["y1"]})

    lines = []
    for row in rows:
        group = sorted(row["words"], key=lambda w: w["x0"])
        heights = [w["y1"] - w["y0"] for w in group]
        line_h = statistics.median(heights)
        # Split purely on line height. Deriving the threshold from the row's own
        # gaps fails on a table row whose cells are all single words ("Alb |
        # Prehipertensi | Normal | ..."): every gap is then a cell gap, so any
        # percentile of them is large and the row stays welded together. A word
        # space runs ~0.3 of line height, while a column gap far exceeds it.
        threshold = line_h * gap_factor

        run = [group[0]]
        for prev, cur in zip(group, group[1:]):
            if cur["x0"] - prev["x1"] > threshold:
                lines.append(make_line(run))
                run = [cur]
            else:
                run.append(cur)
        lines.append(make_line(run))
    return dedupe_lines([ln for ln in lines if ln["text"].strip()])


def make_line(run):
    return {
        "words": list(run),
        "text": " ".join(w["text"] for w in run),
        "conf": sum(w["conf"] for w in run) / len(run),
        "x0": min(w["x0"] for w in run),
        "y0": min(w["y0"] for w in run),
        "x1": max(w["x1"] for w in run),
        "y1": max(w["y1"] for w in run),
    }


def dedupe_lines(lines):
    """Remove lines that restate an overlapping neighbour (polarity doubles)."""
    kept = []
    for ln in sorted(lines, key=lambda ln: -ln["conf"]):
        same = lambda k: ln["text"].strip().lower() == k["text"].strip().lower()
        if any(iou(ln, k) > 0.6 or (same(k) and iou(ln, k) > 0.25) for k in kept):
            continue
        kept.append(ln)
    return kept


# --------------------------------------------------------------------------
# Per-line refinement
# --------------------------------------------------------------------------
TARGET_LINE_PX = 64  # crops are upscaled so glyphs are about this tall


def binarise_ink(gray):
    """Split a crop into ink and background. Returns (ink_mask, ink_is_light).

    Polarity is decided by which class is the MINORITY, because letters always
    cover less of their box than the background does. Judging it from the
    crop's border instead breaks whenever the padded box spills past a coloured
    block -- which sampled the navy "EKG" header as navy-on-white and emitted
    navy text that was invisible against its own cell.
    """
    _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    if (binary > 0).mean() <= 0.5:
        return binary, True  # bright pixels are the minority -> light text
    return 255 - binary, False


def fragmentation(text):
    """Count 1-character alphabetic tokens -- the signature of split words."""
    return sum(1 for t in text.split() if len(t) == 1 and t.isalpha())


def refine_line(rgb, line, lang, trust_conf=88.0, pad_ratio=0.35):
    """Re-OCR an uncertain line in isolation; keep it only if it improves.

    Whole-page segmentation mangles light-on-dark and small text, but psm 7 on
    an upscaled crop has the opposite failure mode -- it sprays spaces inside
    words ("Pa rameter", "Kelaika n"). So this runs only on lines the page pass
    was unsure about, and accepts the rewrite only when confidence rises and
    the text did not fragment.

    Returns (text, confidence).
    """
    if line["conf"] >= trust_conf:
        return line["text"], line["conf"]

    h, w = rgb.shape[:2]
    lh = line["y1"] - line["y0"]
    pad_y, pad_x = max(int(lh * pad_ratio), 4), max(int(lh * 0.4), 6)
    x0, y0 = max(line["x0"] - pad_x, 0), max(line["y0"] - pad_y, 0)
    x1, y1 = min(line["x1"] + pad_x, w), min(line["y1"] + pad_y, h)
    if x1 - x0 < 4 or y1 - y0 < 4:
        return line["text"], line["conf"]

    gray = cv2.cvtColor(rgb[y0:y1, x0:x1], cv2.COLOR_RGB2GRAY)
    if binarise_ink(gray)[1]:  # Tesseract expects dark text on a light ground
        gray = 255 - gray
    scale = TARGET_LINE_PX / max(lh, 1)
    if scale > 1.05:
        gray = cv2.resize(gray, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
    gray = cv2.copyMakeBorder(gray, 12, 12, 12, 12, cv2.BORDER_REPLICATE)

    try:
        data = pytesseract.image_to_data(
            Image.fromarray(gray), lang=lang, config="--psm 7",
            output_type=pytesseract.Output.DICT,
        )
    except pytesseract.TesseractError:
        return line["text"], line["conf"]

    tokens, confs = [], []
    for i in range(len(data["text"])):
        tok = data["text"][i].strip().replace("|", "")
        try:
            c = float(data["conf"][i])
        except (TypeError, ValueError):
            continue
        if tok and c >= 0:
            tokens.append(tok)
            confs.append(c)
    if not tokens:
        return line["text"], line["conf"]

    text, conf = " ".join(tokens), sum(confs) / len(confs)
    if conf < 80.0 or conf <= line["conf"] + 5.0:
        return line["text"], line["conf"]
    if fragmentation(text) > fragmentation(line["text"]):
        return line["text"], line["conf"]
    # A barely-readable original carries no length information worth trusting.
    if line["conf"] > 40.0 and abs(len(text) - len(line["text"])) > max(6, 0.6 * len(line["text"])):
        return line["text"], line["conf"]
    return text, conf


def is_meaningful(text):
    """Reject OCR noise picked up from icons and rules."""
    cleaned = text.replace("|", "").strip()
    return bool(cleaned) and any(ch.isalnum() for ch in cleaned)


def recover_missed_text(rgb, kept, lang, min_conf, dpi, max_candidates=250):
    """Second-chance detection for text the page-layout analysis skipped.

    Whole-page segmentation silently mangles short words sitting on strong
    colour blocks: the navy "Inisial" and "EKG" table headers come back as
    "ed" (conf 18) and "ce" (conf 9), with boxes too poor for refinement to
    anchor to. This locates text-shaped regions that the accepted lines do not
    already cover, and re-reads each one in isolation -- where the same pixels
    read perfectly.
    """
    h, w = rgb.shape[:2]
    heights = [ln["y1"] - ln["y0"] for ln in kept]
    med_h = statistics.median(heights) if heights else dpi / 10.0

    covered = np.zeros((h, w), dtype=np.uint8)
    for ln in kept:
        covered[max(ln["y0"], 0):ln["y1"], max(ln["x0"], 0):ln["x1"]] = 1

    gray = cv2.cvtColor(rgb, cv2.COLOR_RGB2GRAY)
    grad = cv2.morphologyEx(gray, cv2.MORPH_GRADIENT, np.ones((3, 3), np.uint8))
    _, binary = cv2.threshold(grad, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    kernel = np.ones((max(1, int(med_h * 0.2)), max(3, int(med_h * 0.9))), np.uint8)
    joined = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)
    contours, _ = cv2.findContours(joined, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    boxes = []
    for contour in contours:
        x, y, cw, ch = cv2.boundingRect(contour)
        if not (0.45 * med_h <= ch <= 3.0 * med_h) or cw < 0.45 * med_h or cw > 40 * ch:
            continue
        if covered[y:y + ch, x:x + cw].mean() > 0.25:
            continue
        boxes.append((cw * ch, {"text": "", "conf": 0.0,
                                "x0": x, "y0": y, "x1": x + cw, "y1": y + ch}))

    # Speculative regions must clear a high bar. Line-art icons OCR into short
    # plausible strings ("AN", "NA") at ordinary confidence, and accepting one
    # is doubly costly: it invents text AND erases the icon underneath it.
    floor = max(min_conf, 80.0)
    found = []
    for _, cand in sorted(boxes, key=lambda b: -b[0])[:max_candidates]:
        text, conf = refine_line(rgb, cand, lang)
        if conf < floor or not is_meaningful(text):
            continue
        cand["text"], cand["conf"] = text.replace("|", "").strip(), conf
        if any(iou(cand, other) > 0.3 or containment(cand, other) > 0.5
               for other in kept + found):
            continue
        found.append(cand)
    return found


def recover_rows(rgb, existing, lang, min_conf):
    """Re-read each text row as an isolated full-width strip.

    A crop holding a single row of text is segmented far more reliably than a
    whole page: the navy table header that the page pass returns as "ed" and
    "ce" comes back complete and correct from its own strip. Contour recovery
    cannot reach it, because the header's solid colour block forms one large
    contour that swallows the words sitting inside it.
    """
    if not existing:
        return []

    bands = []
    for ln in sorted(existing, key=lambda l: l["y0"]):
        for band in bands:
            if same_row(ln, band):
                band["y0"] = min(band["y0"], ln["y0"])
                band["y1"] = max(band["y1"], ln["y1"])
                break
        else:
            bands.append({"y0": ln["y0"], "y1": ln["y1"]})

    # Recovered text is speculative, so it must clear a much higher bar than
    # the primary pass: strip re-OCR reports table rules as "|", re-reads
    # slivers of captured words, and turns line-art icons into short plausible
    # strings ("AN", "NA" at conf 73-84). Genuine finds score far higher --
    # the missing "Inisial" and "EKG" headers come back at 95 and 91.
    floor = max(min_conf, 85.0)
    h, w = rgb.shape[:2]
    words = []
    for band in bands:
        bh = band["y1"] - band["y0"]
        if bh < 6:
            continue
        pad = max(int(bh * 0.3), 4)
        y0, y1 = max(band["y0"] - pad, 0), min(band["y1"] + pad, h)
        strip = cv2.cvtColor(rgb[y0:y1, :], cv2.COLOR_RGB2GRAY)
        scale = TARGET_LINE_PX / float(bh)
        if scale > 1.05:
            strip = cv2.resize(strip, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
        else:
            scale = 1.0

        for variant in (strip, 255 - strip):
            try:
                data = pytesseract.image_to_data(
                    Image.fromarray(variant), lang=lang, config="--psm 6",
                    output_type=pytesseract.Output.DICT,
                )
            except pytesseract.TesseractError:
                continue
            for i in range(len(data["text"])):
                text = data["text"][i].replace("|", "").strip()
                try:
                    conf = float(data["conf"][i])
                except (TypeError, ValueError):
                    continue
                if not text or conf < floor:
                    continue
                words.append(
                    {
                        "text": text, "conf": conf,
                        "x0": int(data["left"][i] / scale),
                        "y0": int(data["top"][i] / scale) + y0,
                        "x1": int((data["left"][i] + data["width"][i]) / scale),
                        "y1": int((data["top"][i] + data["height"][i]) / scale) + y0,
                    }
                )

    novel = lambda box: not any(
        iou(box, e) > 0.2 or containment(box, e) > 0.5 for e in existing
    )
    fresh = [c for c in dedupe_boxes(words) if novel(c)]
    return [
        ln for ln in assemble_lines(fresh)
        if is_meaningful(ln["text"]) and len(ln["text"].strip()) >= 2 and novel(ln)
    ]


# --------------------------------------------------------------------------
# Pixel-level styling
# --------------------------------------------------------------------------
def glyph_mask(rgb, line, pad=2):
    """Ink mask for a line, plus the colour of the surface it sits on.

    Returns (mask, box, background_colour, background_share).

    The background is measured from a ring just outside the line, and ink is
    then whatever differs from it. This beats deciding polarity from the crop
    alone: "EKG" is three bold capitals in a tight box, so its white letters
    are not the minority of their own pixels, and a minority-class rule painted
    them navy -- invisible against the navy header they sit on.
    """
    h, w = rgb.shape[:2]
    line_h = max(line["y1"] - line["y0"], 1)
    extra = max(3, int(line_h * 0.15))

    ox0, oy0 = max(line["x0"] - pad - extra, 0), max(line["y0"] - pad - extra, 0)
    ox1, oy1 = min(line["x1"] + pad + extra, w), min(line["y1"] + pad + extra, h)
    x0, y0 = max(line["x0"] - pad, 0), max(line["y0"] - pad, 0)
    x1, y1 = min(line["x1"] + pad, w), min(line["y1"] + pad, h)
    if x1 <= x0 or y1 <= y0 or ox1 <= ox0 or oy1 <= oy0:
        return None, (x0, y0, x1, y1), np.array([255, 255, 255], np.uint8), 0.0

    outer = rgb[oy0:oy1, ox0:ox1]
    thickness = max(2, extra // 2)
    ring = np.concatenate([
        outer[:thickness].reshape(-1, 3), outer[-thickness:].reshape(-1, 3),
        outer[:, :thickness].reshape(-1, 3), outer[:, -thickness:].reshape(-1, 3),
    ])
    if len(ring) < 10:
        return None, (x0, y0, x1, y1), np.array([255, 255, 255], np.uint8), 0.0
    background, share = dominant_background(ring)

    region = rgb[y0:y1, x0:x1]
    distance = np.abs(region.astype(np.int16) - background.astype(np.int16)).max(axis=2)
    _, mask = cv2.threshold(distance.astype(np.uint8), 0, 255,
                            cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return mask, (x0, y0, x1, y1), background, share


def measure_ink(rgb, box):
    """Colour, boldness and ink height of the glyphs inside `box`."""
    mask, (x0, y0, x1, y1), _, _ = glyph_mask(rgb, box)
    if mask is None or not mask.any():
        return None
    ink = mask > 0
    # Sample colour from the glyph CORE: a letter's outer pixels are the
    # anti-aliased edge plus this deck's soft glow, which drags the median
    # toward the glow's hue.
    core = cv2.erode(mask, np.ones((3, 3), np.uint8), iterations=1) > 0
    sample = core if core.sum() >= 20 else ink
    pixels = rgb[y0:y1, x0:x1][sample]
    colour = tuple(int(v) for v in np.median(pixels, axis=0)) if len(pixels) else (0, 0, 0)

    rows = np.where(ink.any(axis=1))[0]
    ink_h = (rows[-1] - rows[0] + 1) if len(rows) else (y1 - y0)
    stroke = float(np.percentile(cv2.distanceTransform(mask, cv2.DIST_L2, 3)[ink], 80)) * 2.0
    return colour, (stroke / max(ink_h, 1)) > 0.115, ink_h


def word_runs(rgb, line, colour_tolerance=40):
    """Split a line into runs of consistent colour and weight.

    This deck emphasises phrases mid-sentence in gold bold. Styling a whole
    line uniformly repaints those sentences entirely gold or entirely navy, so
    each word is measured separately and neighbouring words of matching style
    are merged back into one run.
    """
    words = line.get("words")
    if not words or " ".join(w["text"] for w in words) != line["text"]:
        return None  # refinement rewrote the text; word boxes no longer align

    groups = []
    for word in words:
        measured = measure_ink(rgb, word)
        if measured is None:
            return None
        colour = measured[0]
        if groups and sum(abs(a - b) for a, b in zip(colour, groups[-1]["color"])) < colour_tolerance:
            groups[-1]["words"].append(word)
        else:
            groups.append({"words": [word], "color": colour})

    # Weight is measured per RUN, not per word. Stroke-to-height ratio swings
    # with which letters a word happens to contain, so measuring each word
    # separately made single sentences alternate bold and regular.
    runs = []
    for group in groups:
        span = {
            "x0": min(w["x0"] for w in group["words"]),
            "y0": min(w["y0"] for w in group["words"]),
            "x1": max(w["x1"] for w in group["words"]),
            "y1": max(w["y1"] for w in group["words"]),
        }
        measured = measure_ink(rgb, span)
        colour, bold = (measured[0], measured[1]) if measured else (group["color"], False)
        runs.append({
            "text": " ".join(w["text"] for w in group["words"]),
            "color": colour, "bold": bold,
        })
    return runs if len(runs) > 1 else None


_FONT_CACHE = {}


def fit_size_to_width(text, width_in, font_path, probe_px=100):
    """Point size at which `text` in `font_path` spans `width_in` inches."""
    if not text or not font_path or width_in <= 0:
        return None
    try:
        font = _FONT_CACHE.get((font_path, probe_px))
        if font is None:
            font = ImageFont.truetype(font_path, probe_px)
            _FONT_CACHE[(font_path, probe_px)] = font
        measured = font.getlength(text)
    except Exception:
        return None
    return 72.0 * width_in * probe_px / measured if measured > 0 else None


def line_style(rgb, line, scale, font_path, bold_font_path):
    """Measure colour, weight and point size from the line's own pixels.

    `scale` is inches-of-slide per pixel. Sizes must be derived in slide space,
    not render space: a 400-DPI render of this deck is 19.1in wide but lands on
    a 13.3in slide, so measuring in render inches oversized every font by 1.43x.
    """
    style = {"color": (0, 0, 0), "bold": False, "size_pt": 12.0, "runs": None}
    measured = measure_ink(rgb, line)
    if measured is None:
        return style
    style["color"], style["bold"], ink_h = measured

    x0, y0, x1, y1 = (max(line["x0"], 0), max(line["y0"], 0), line["x1"], line["y1"])
    style["runs"] = word_runs(rgb, line)

    size_from_height = (ink_h * scale * 72.0) / 0.95  # ascender-to-descender ~= 0.95 em
    width_in = (x1 - x0) * scale
    size_from_width = fit_size_to_width(
        line["text"], width_in, bold_font_path if style["bold"] else font_path
    )
    if size_from_width:
        lo, hi = size_from_height * 0.75, size_from_height * 1.25
        style["size_pt"] = max(lo, min(hi, size_from_width))
    else:
        style["size_pt"] = size_from_height
    style["size_pt"] = max(5.0, min(120.0, style["size_pt"]))
    return style


# --------------------------------------------------------------------------
# Erasing captured text from the page
# --------------------------------------------------------------------------
UNIFORM_SHARE = 0.70  # fraction of background pixels that must match the dominant colour


def dominant_background(pixels, tolerance=14):
    """Most common colour among `pixels`, plus the share of pixels matching it.

    Uses the mode rather than the mean or median because this deck sets gold
    and blue text in a soft glow. A mean or median is dragged into the glow and
    fills the box with a pale tint; the glow's pixels are spread across many
    near-but-distinct values, so the true background still wins the mode.
    """
    quantised = (pixels // 8).astype(np.int32)
    packed = quantised[:, 0] * 4096 + quantised[:, 1] * 64 + quantised[:, 2]
    values, counts = np.unique(packed, return_counts=True)
    top = values[counts.argmax()]
    centre = np.array([(top // 4096) % 64, (top // 64) % 64, top % 64]) * 8 + 4
    matching = np.abs(pixels.astype(np.int16) - centre).max(axis=1) <= tolerance
    # Average the matching pixels rather than returning the bucket centre --
    # the centre of white's bucket is 252, which paints a faintly grey block
    # onto a pure-white page.
    colour = np.median(pixels[matching], axis=0).astype(np.uint8)
    return colour, float(matching.mean())


def erase_text(rgb, lines, dilate=None):
    """Remove captured text from the page image.

    Two strategies, chosen per line:
      * uniform background (a table cell, a navy header, plain white) -> repaint
        the bbox with the measured background colour. Exact and fast.
      * textured background (gradients, charts, photos) -> inpaint the dilated
        glyph mask, reconstructing what was behind the letters.

    v1 always did a flat fill, using the MEAN of a border strip that crossed
    table rules and neighbouring coloured cells -- a colour present nowhere on
    the page, which is where the yellow and blue blotches came from. Sampling
    the actual background pixels inside the bbox, and only when they are
    genuinely uniform, removes that failure mode.
    """
    out = rgb.copy()
    h, w = rgb.shape[:2]
    inpaint_mask = np.zeros((h, w), dtype=np.uint8)

    for line in lines:
        line_h = line["y1"] - line["y0"]
        pad = max(3, int(line_h * 0.16))  # wide enough to swallow the glow
        mask, (x0, y0, x1, y1), colour, share = glyph_mask(rgb, line, pad=pad)
        if mask is None or not mask.any():
            continue

        if share > UNIFORM_SHARE:
            out[y0:y1, x0:x1] = colour
        else:
            grown = cv2.dilate(
                mask, np.ones((dilate or max(7, int(line_h * 0.18)),) * 2, np.uint8),
                iterations=1,
            )
            inpaint_mask[y0:y1, x0:x1] = np.maximum(inpaint_mask[y0:y1, x0:x1], grown)

    if inpaint_mask.any():
        bgr = cv2.cvtColor(out, cv2.COLOR_RGB2BGR)
        out = cv2.cvtColor(cv2.inpaint(bgr, inpaint_mask, 4, cv2.INPAINT_TELEA), cv2.COLOR_BGR2RGB)
    return out


# --------------------------------------------------------------------------
# Paragraph grouping
# --------------------------------------------------------------------------
def group_paragraphs(items, gap_factor=0.55):
    """Merge stacked lines that clearly belong to one block.

    Deliberately conservative -- lines must share a column, sit close together,
    and match in size and colour -- so table rows stay separate.
    """
    paragraphs = []
    for it in sorted(items, key=lambda it: (it["line"]["y0"], it["line"]["x0"])):
        line, style = it["line"], it["style"]
        merged = False
        if paragraphs:
            cur = paragraphs[-1]
            last = cur["items"][-1]
            lh = last["line"]["y1"] - last["line"]["y0"]
            gap = line["y0"] - last["line"]["y1"]
            overlap = min(line["x1"], cur["x1"]) - max(line["x0"], cur["x0"])
            span = min(line["x1"] - line["x0"], cur["x1"] - cur["x0"])
            aligned = (
                abs(line["x0"] - cur["x0"]) < lh * 1.2
                or abs((line["x0"] + line["x1"]) / 2 - (cur["x0"] + cur["x1"]) / 2) < lh * 1.2
            )
            if (
                -lh * 0.25 <= gap <= lh * gap_factor
                and 0.85 <= style["size_pt"] / max(last["style"]["size_pt"], 0.1) <= 1.18
                and sum(abs(a - b) for a, b in zip(style["color"], last["style"]["color"])) < 90
                and overlap > 0.55 * max(span, 1)
                and aligned
            ):
                cur["items"].append(it)
                cur["x0"], cur["y0"] = min(cur["x0"], line["x0"]), min(cur["y0"], line["y0"])
                cur["x1"], cur["y1"] = max(cur["x1"], line["x1"]), max(cur["y1"], line["y1"])
                merged = True
        if not merged:
            paragraphs.append(
                {"items": [it], "x0": line["x0"], "y0": line["y0"],
                 "x1": line["x1"], "y1": line["y1"]}
            )
    return paragraphs


def detect_alignment(para):
    if len(para["items"]) < 2:
        return PP_ALIGN.LEFT
    lefts = [it["line"]["x0"] for it in para["items"]]
    rights = [it["line"]["x1"] for it in para["items"]]
    centers = [(l + r) / 2 for l, r in zip(lefts, rights)]
    spread = lambda v: max(v) - min(v)
    options = {
        PP_ALIGN.LEFT: spread(lefts),
        PP_ALIGN.RIGHT: spread(rights),
        PP_ALIGN.CENTER: spread(centers),
    }
    return min(options, key=options.get)


# --------------------------------------------------------------------------
# Build
# --------------------------------------------------------------------------
def build_pptx(pdf_path, out_path, lang="ind+eng", dpi=400, min_conf=35,
               background="clean", font_name="Arial", font_path=None,
               bold_font_path=None, refine=True, recover=True, bg_max_width=2400,
               pages=None):
    doc = pymupdf.open(pdf_path)
    if len(doc) == 0:
        raise ValueError("PDF has no pages")
    selected = pages if pages else range(len(doc))

    rect = doc[0].rect
    slide_w_in = 13.333
    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_w_in * (rect.height / rect.width))
    blank = prs.slide_layouts[6]

    for pno in selected:
        rgb = render_page(doc[pno], dpi)
        scale = slide_w_in / rgb.shape[1]  # inches per pixel

        lines = assemble_lines(detect_words(rgb, lang))

        kept = []
        for line in lines:
            text, conf = refine_line(rgb, line, lang) if refine else (line["text"], line["conf"])
            # Confidence is judged only after refinement, so text the page pass
            # merely stumbled over is repaired rather than discarded.
            if conf < min_conf or not is_meaningful(text):
                continue
            line["text"] = text.replace("|", "").strip()
            line["conf"] = conf
            kept.append(line)

        recovered = []
        if recover:
            recovered += recover_missed_text(rgb, kept, lang, min_conf, dpi)
            recovered += recover_rows(rgb, kept + recovered, lang, min_conf)
        kept.extend(recovered)

        items = [
            {"line": ln, "style": line_style(rgb, ln, scale, font_path, bold_font_path)}
            for ln in kept
        ]
        paragraphs = group_paragraphs(items)

        slide = prs.slides.add_slide(blank)
        if background != "none":
            # Only `kept` lines are erased: anything unreadable stays as pixels.
            plate = erase_text(rgb, kept) if background == "clean" else rgb
            slide.shapes.add_picture(
                encode_plate(plate, bg_max_width), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
        for para in paragraphs:
            add_text_block(slide, para, scale, font_name)

        print(
            f"  slide {pno + 1}/{len(doc)}: {len(lines)} detected, {len(kept)} kept "
            f"(+{len(recovered)} recovered) -> {len(paragraphs)} text box(es)",
            file=sys.stderr,
        )

    prs.save(out_path)


def encode_plate(plate, max_width):
    """Downscale the background before embedding.

    OCR needs a heavily upscaled page, but storing that upscale is pure bloat --
    it carries no detail the source did not have, and it pushed the deck to
    34 MB. Downscaling to a sane width keeps the visual result identical.
    """
    img = Image.fromarray(plate)
    if max_width and img.width > max_width:
        height = round(img.height * max_width / img.width)
        img = img.resize((max_width, height), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    buf.seek(0)
    return buf


def add_text_block(slide, para, scale, font_name):
    pad_in = 0.04  # breathing room so PowerPoint never re-wraps a fitted line
    box = slide.shapes.add_textbox(
        Inches(max(para["x0"] * scale - pad_in, 0)),
        Inches(max(para["y0"] * scale - pad_in * 0.5, 0)),
        Inches((para["x1"] - para["x0"]) * scale + pad_in * 2),
        Inches((para["y1"] - para["y0"]) * scale + pad_in),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    align = detect_alignment(para)
    for i, it in enumerate(para["items"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        style = it["style"]
        size = Pt(round(style["size_pt"], 1))
        pieces = style.get("runs") or [
            {"text": it["line"]["text"], "color": style["color"], "bold": style["bold"]}
        ]
        for index, piece in enumerate(pieces):
            run = p.add_run()
            run.text = piece["text"] if index == 0 else " " + piece["text"]
            run.font.size = size
            run.font.bold = piece["bold"]
            run.font.name = font_name
            run.font.color.rgb = RGBColor(*piece["color"])


def parse_pages(spec):
    """Turn "3", "1-3" or "1,4,7" into a list of 0-based page indices."""
    if not spec:
        return None
    pages = []
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            first, last = (int(v) for v in part.split("-", 1))
            pages.extend(range(first - 1, last))
        elif part:
            pages.append(int(part) - 1)
    return pages


def main():
    ap = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter
    )
    ap.add_argument("pdf_path")
    ap.add_argument("out_path")
    ap.add_argument("--lang", default="ind+eng", help="Tesseract language(s) (default: ind+eng)")
    ap.add_argument("--dpi", type=int, default=400, help="OCR render resolution (default: 400)")
    ap.add_argument("--min-conf", type=int, default=35,
                    help="Drop text still below this confidence after refinement")
    ap.add_argument("--background", choices=["clean", "original", "none"], default="clean",
                    help="clean = captured text erased (default); original = untouched; none = blank")
    ap.add_argument("--bg-max-width", type=int, default=2400,
                    help="Downscale the background to this width before embedding")
    ap.add_argument("--font", default="Arial", help="Font for the generated text boxes")
    ap.add_argument("--font-file", default=r"C:\Windows\Fonts\arial.ttf")
    ap.add_argument("--bold-font-file", default=r"C:\Windows\Fonts\arialbd.ttf")
    ap.add_argument("--no-refine", action="store_true", help="Skip the per-line re-OCR pass")
    ap.add_argument("--no-recover", action="store_true",
                    help="Skip the second-chance pass for text the layout analysis missed")
    ap.add_argument("--tesseract", help="Full path to tesseract.exe")
    ap.add_argument("--pages", help='Pages to convert, 1-based, e.g. "3" or "1-3" or "1,4,7"')
    args = ap.parse_args()

    pytesseract.pytesseract.tesseract_cmd = locate_tesseract(args.tesseract)
    font_file = args.font_file if os.path.isfile(args.font_file) else None
    bold_file = args.bold_font_file if os.path.isfile(args.bold_font_file) else font_file

    build_pptx(
        args.pdf_path, args.out_path,
        lang=args.lang, dpi=args.dpi, min_conf=args.min_conf,
        background=args.background, font_name=args.font,
        font_path=font_file, bold_font_path=bold_file,
        refine=not args.no_refine, recover=not args.no_recover,
        bg_max_width=args.bg_max_width, pages=parse_pages(args.pages),
    )
    print(f"Saved: {args.out_path}", file=sys.stderr)


if __name__ == "__main__":
    main()
